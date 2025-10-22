"""
Çoklu Dosya Formatı İşleyici
PDF, DOCX, TXT, PPTX dosyalarını text'e çevirir
"""

import os
from typing import List, Dict
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation

class DocumentProcessor:
    """Farklı dosya formatlarını işleyen sınıf"""
    
    def __init__(self):
        self.supported_formats = ['.pdf', '.docx', '.txt', '.pptx']
    
    def extract_text(self, file_path: str) -> Dict[str, any]:
        """
        Dosyadan metin çıkarır
        
        Args:
            file_path: Dosya yolu
            
        Returns:
            dict: {
                'text': str,
                'metadata': dict,
                'page_count': int
            }
        """
        file_extension = os.path.splitext(file_path)[1].lower()
        
        if file_extension not in self.supported_formats:
            raise ValueError(f"Desteklenmeyen format: {file_extension}")
        
        if file_extension == '.pdf':
            return self._process_pdf(file_path)
        elif file_extension == '.docx':
            return self._process_docx(file_path)
        elif file_extension == '.txt':
            return self._process_txt(file_path)
        elif file_extension == '.pptx':
            return self._process_pptx(file_path)
    
    def _process_pdf(self, file_path: str) -> Dict:
        """PDF dosyasını işler"""
        try:
            reader = PdfReader(file_path)
            text = ""
            
            for page in reader.pages:
                text += page.extract_text() + "\n\n"
            
            return {
                'text': text,
                'metadata': {
                    'file_type': 'pdf',
                    'file_name': os.path.basename(file_path)
                },
                'page_count': len(reader.pages)
            }
        except Exception as e:
            raise Exception(f"PDF işleme hatası: {str(e)}")
    
    def _process_docx(self, file_path: str) -> Dict:
        """DOCX dosyasını işler"""
        try:
            doc = Document(file_path)
            text = ""
            
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            
            return {
                'text': text,
                'metadata': {
                    'file_type': 'docx',
                    'file_name': os.path.basename(file_path)
                },
                'page_count': len(doc.paragraphs)
            }
        except Exception as e:
            raise Exception(f"DOCX işleme hatası: {str(e)}")
    
    def _process_txt(self, file_path: str) -> Dict:
        """TXT dosyasını işler"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                text = f.read()
            
            return {
                'text': text,
                'metadata': {
                    'file_type': 'txt',
                    'file_name': os.path.basename(file_path)
                },
                'page_count': 1
            }
        except Exception as e:
            raise Exception(f"TXT işleme hatası: {str(e)}")
    
    def _process_pptx(self, file_path: str) -> Dict:
        """PPTX dosyasını işler"""
        try:
            prs = Presentation(file_path)
            text = ""
            
            for slide_num, slide in enumerate(prs.slides, 1):
                text += f"\n--- Slayt {slide_num} ---\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            
            return {
                'text': text,
                'metadata': {
                    'file_type': 'pptx',
                    'file_name': os.path.basename(file_path)
                },
                'page_count': len(prs.slides)
            }
        except Exception as e:
            raise Exception(f"PPTX işleme hatası: {str(e)}")
    
    def chunk_text(self, text: str, chunk_size: int = 1000, overlap: int = 200) -> List[str]:
        """
        Metni daha küçük parçalara böler
        
        Args:
            text: İşlenecek metin
            chunk_size: Her parçanın boyutu
            overlap: Parçalar arası örtüşme
            
        Returns:
            list: Metin parçaları
        """
        chunks = []
        start = 0
        text_length = len(text)
        
        while start < text_length:
            end = start + chunk_size
            chunk = text[start:end]
            chunks.append(chunk)
            start += chunk_size - overlap
        
        return chunks
